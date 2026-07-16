"""
tools/tools_media.py — Document intelligence, image/video analysis, facial recognition.

Capabilities:
  read_document(path)          — extract text+tables from PDF, DOCX, PPTX, TXT
  analyze_image(path, q)       — Claude vision analysis of any image
  analyze_video(path, q)       — sample video frames → Claude vision
  detect_faces(path)           — DeepFace: detect faces + age/gender/emotion
  recognize_face(path)         — match face against data/faces/ database
  register_face(path, name)    — add image to face database
  ocr_image(path)              — extract text from image via Tesseract / Claude fallback

Face database: data/faces/<Name>/img1.jpg  (gitignored — personal biometric data)
"""

import base64
import io
import os
import re
import shutil
from pathlib import Path
from typing import Dict, List, Optional, Tuple

from agent.observability import track_silent

_ROOT       = Path(__file__).parent.parent
_FACES_DIR  = _ROOT / "data" / "faces"
_TEMP_DIR   = _ROOT / "data" / "media_temp"

# ── optional imports ──────────────────────────────────────────────────────────

try:
    import pdfplumber
    _PDF_OK = True
except ImportError:
    _PDF_OK = False

try:
    from docx import Document as _DocxDocument
    _DOCX_OK = True
except ImportError:
    _DOCX_OK = False

try:
    from pptx import Presentation as _Pptx
    _PPTX_OK = True
except ImportError:
    _PPTX_OK = False

try:
    from PIL import Image as _PIL
    _PIL_OK = True
except ImportError:
    _PIL_OK = False

try:
    import cv2 as _cv2
    _CV2_OK = True
except ImportError:
    _CV2_OK = False

_DEEPFACE_OK = None   # None=untried  True=ok  False=missing
_deepface_ref = None

def _get_deepface():
    global _DEEPFACE_OK, _deepface_ref
    if _DEEPFACE_OK is None:
        try:
            from deepface import DeepFace as _df
            _deepface_ref = _df
            _DEEPFACE_OK = True
        except ImportError:
            _DEEPFACE_OK = False
    return _deepface_ref

try:
    import pytesseract
    _TESS_CMD = os.environ.get("TESSERACT_CMD", "")
    if _TESS_CMD and os.path.exists(_TESS_CMD):
        pytesseract.pytesseract.tesseract_cmd = _TESS_CMD
    _TESS_OK = True
except ImportError:
    _TESS_OK = False

try:
    import anthropic as _anthropic
    _ANTHROPIC_OK = True
except ImportError:
    _ANTHROPIC_OK = False

try:
    from google import genai as _genai
    _GENAI_OK = True
except ImportError:
    _GENAI_OK = False


# ---------------------------------------------------------------------------
# Internal helpers
# ---------------------------------------------------------------------------

def _encode_image(path: str) -> Tuple[str, str]:
    """Return (base64_data, media_type) for an image file."""
    ext = Path(path).suffix.lower()
    mime_map = {
        ".jpg":  "image/jpeg",
        ".jpeg": "image/jpeg",
        ".png":  "image/png",
        ".gif":  "image/gif",
        ".webp": "image/webp",
    }
    media_type = mime_map.get(ext, "image/jpeg")
    with open(path, "rb") as f:
        data = base64.standard_b64encode(f.read()).decode("utf-8")
    return data, media_type


def _claude_vision(image_paths: List[str], question: str, max_tokens: int = 2048) -> str:
    """Send one or more images to Claude vision and return the response text."""
    if not _ANTHROPIC_OK:
        return "[vision] anthropic not installed"
    try:
        from dotenv import load_dotenv
        load_dotenv(_ROOT / ".env")
    except Exception:
        pass

    key = os.environ.get("ANTHROPIC_API_KEY")
    if not key:
        return "[vision] ANTHROPIC_API_KEY not set"

    client = _anthropic.Anthropic(api_key=key)
    content = []
    for path in image_paths[:5]:  # Claude supports up to 20, but keep reasonable
        try:
            data, mime = _encode_image(path)
            content.append({
                "type": "image",
                "source": {"type": "base64", "media_type": mime, "data": data},
            })
        except Exception as e:
            content.append({"type": "text", "text": f"[Could not load {path}: {e}]"})

    content.append({"type": "text", "text": question or "Describe this image in detail."})

    msg = client.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=max_tokens,
        messages=[{"role": "user", "content": content}],
    )
    return msg.content[0].text


def _gemini_vision(image_paths: List[str], question: str) -> str:
    """Send images to Gemini Vision (free tier) and return the response text.

    Legacy permissive wrapper — swallows errors and returns a sentinel string.
    For T-121 provider-chain dispatch, use _gemini_vision_strict via the router.
    """
    if not _GENAI_OK:
        return "[gemini_vision] google-genai not installed"
    try:
        from dotenv import load_dotenv
        load_dotenv(_ROOT / ".env")
    except Exception:
        pass

    key = os.environ.get("GEMINI_API_KEY") or os.environ.get("GOOGLE_API_KEY")
    if not key:
        return "[gemini_vision] GEMINI_API_KEY not set"

    try:
        client = _genai.Client(api_key=key)
        parts = []
        for path in image_paths[:5]:
            data, mime = _encode_image(path)
            import base64 as _b64
            raw = _b64.b64decode(data)
            parts.append({"inline_data": {"mime_type": mime, "data": raw}})
        parts.append(question or "Describe this image in detail.")
        response = client.models.generate_content(
            model="gemini-2.0-flash",
            contents=parts,
        )
        return response.text
    except Exception as e:
        return f"[gemini_vision error] {e}"


# ── T-121 strict provider wrappers — raise instead of returning sentinel ──────
# These are the callables registered with the ProviderRouter. They re-raise
# any exception (so the router can detect 429s and advance), and treat
# "unavailable" preconditions (missing lib/key) as hard skips by raising
# ProviderError so the router moves on without recording it as a 429.

def _gemini_flash_strict(image_paths: List[str], question: str, max_tokens: int = 2048) -> str:
    """Gemini 2.0 Flash. Raises on any error. Router converts 429-shaped errors."""
    from agent.provider_router import ProviderError
    if not _GENAI_OK:
        raise ProviderError("google-genai not installed")
    try:
        from dotenv import load_dotenv
        load_dotenv(_ROOT / ".env")
    except Exception:
        pass
    key = os.environ.get("GEMINI_API_KEY") or os.environ.get("GOOGLE_API_KEY")
    if not key:
        raise ProviderError("GEMINI_API_KEY not set")

    client = _genai.Client(api_key=key)
    parts = []
    for path in image_paths[:5]:
        data, mime = _encode_image(path)
        import base64 as _b64
        raw = _b64.b64decode(data)
        parts.append({"inline_data": {"mime_type": mime, "data": raw}})
    parts.append(question or "Describe this image in detail.")
    response = client.models.generate_content(model="gemini-2.0-flash", contents=parts)
    return response.text


def _gemini_pro_strict(image_paths: List[str], question: str, max_tokens: int = 2048) -> str:
    """Gemini 1.5 Pro. Higher quota than Flash; cheaper than Claude. Raises on any error."""
    from agent.provider_router import ProviderError
    if not _GENAI_OK:
        raise ProviderError("google-genai not installed")
    try:
        from dotenv import load_dotenv
        load_dotenv(_ROOT / ".env")
    except Exception:
        pass
    key = os.environ.get("GEMINI_API_KEY") or os.environ.get("GOOGLE_API_KEY")
    if not key:
        raise ProviderError("GEMINI_API_KEY not set")

    client = _genai.Client(api_key=key)
    parts = []
    for path in image_paths[:5]:
        data, mime = _encode_image(path)
        import base64 as _b64
        raw = _b64.b64decode(data)
        parts.append({"inline_data": {"mime_type": mime, "data": raw}})
    parts.append(question or "Describe this image in detail.")
    response = client.models.generate_content(model="gemini-1.5-pro", contents=parts)
    return response.text


def _claude_haiku_vision_strict(image_paths: List[str], question: str, max_tokens: int = 2048) -> str:
    """Claude Haiku 4.5 vision — cheap paid fallback. Raises on any error."""
    from agent.provider_router import ProviderError
    if not _ANTHROPIC_OK:
        raise ProviderError("anthropic not installed")
    try:
        from dotenv import load_dotenv
        load_dotenv(_ROOT / ".env")
    except Exception:
        pass
    key = os.environ.get("ANTHROPIC_API_KEY")
    if not key:
        raise ProviderError("ANTHROPIC_API_KEY not set")

    client = _anthropic.Anthropic(api_key=key)
    content = []
    for path in image_paths[:5]:
        data, mime = _encode_image(path)
        content.append({"type": "image", "source": {"type": "base64", "media_type": mime, "data": data}})
    content.append({"type": "text", "text": question or "Describe this image in detail."})
    msg = client.messages.create(
        model="claude-haiku-4-5-20251001",
        max_tokens=max_tokens,
        messages=[{"role": "user", "content": content}],
    )
    return msg.content[0].text


def _claude_sonnet_vision_strict(image_paths: List[str], question: str, max_tokens: int = 2048) -> str:
    """Claude Sonnet 4.6 vision — last-resort paid fallback. Raises on any error."""
    from agent.provider_router import ProviderError
    if not _ANTHROPIC_OK:
        raise ProviderError("anthropic not installed")
    try:
        from dotenv import load_dotenv
        load_dotenv(_ROOT / ".env")
    except Exception:
        pass
    key = os.environ.get("ANTHROPIC_API_KEY")
    if not key:
        raise ProviderError("ANTHROPIC_API_KEY not set")

    client = _anthropic.Anthropic(api_key=key)
    content = []
    for path in image_paths[:5]:
        data, mime = _encode_image(path)
        content.append({"type": "image", "source": {"type": "base64", "media_type": mime, "data": data}})
    content.append({"type": "text", "text": question or "Describe this image in detail."})
    msg = client.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=max_tokens,
        messages=[{"role": "user", "content": content}],
    )
    return msg.content[0].text


# Singleton router built on first use
_VISION_ROUTER = None

def _get_vision_router():
    """Build (or return cached) ProviderRouter for vision calls."""
    global _VISION_ROUTER
    if _VISION_ROUTER is None:
        from agent.provider_router import ProviderRouter
        _VISION_ROUTER = ProviderRouter(
            name="vision",
            providers=[
                ("gemini_flash", _gemini_flash_strict),
                ("gemini_pro", _gemini_pro_strict),
                ("claude_haiku", _claude_haiku_vision_strict),
                ("claude_sonnet", _claude_sonnet_vision_strict),
            ],
        )
    return _VISION_ROUTER


def _vision_via_router(image_paths: List[str], question: str, max_tokens: int = 2048) -> Tuple[str, str]:
    """Dispatch a vision call through the router. Returns (text, provider_used).

    Raises AllProvidersExhausted if every provider failed/was cooled.
    """
    router = _get_vision_router()
    # The router doesn't tell us which provider succeeded — wrap to capture it
    successful_provider = {"name": "unknown"}
    original_providers = list(router.providers)

    def _wrap(name, fn):
        def _wrapped(*a, **kw):
            r = fn(*a, **kw)
            successful_provider["name"] = name
            return r
        return _wrapped

    router.providers = [(n, _wrap(n, fn)) for n, fn in original_providers]
    try:
        result = router.call(image_paths, question, max_tokens=max_tokens)
        return result, successful_provider["name"]
    finally:
        router.providers = original_providers


def _ensure_dir(d: Path):
    d.mkdir(parents=True, exist_ok=True)


# ---------------------------------------------------------------------------
# MediaTools
# ---------------------------------------------------------------------------

class MediaTools:
    """Document intelligence + image/video analysis + facial recognition."""

    # ── Document reading ───────────────────────────────────────────────────────

    @staticmethod
    def read_document(path: str, max_chars: int = 50_000) -> Dict:
        """
        Extract text and tables from PDF, DOCX, PPTX, or TXT.
        Returns {"success", "text", "pages", "tables", "metadata"}.
        """
        path = str(path)
        ext  = Path(path).suffix.lower()

        if not Path(path).exists():
            return {"success": False, "error": f"File not found: {path}"}

        try:
            if ext == ".pdf":
                return MediaTools._read_pdf(path, max_chars)
            elif ext in (".docx", ".doc"):
                return MediaTools._read_docx(path, max_chars)
            elif ext in (".pptx", ".ppt"):
                return MediaTools._read_pptx(path, max_chars)
            elif ext in (".txt", ".md", ".csv", ".json", ".py", ".js", ".ts"):
                text = Path(path).read_text(encoding="utf-8", errors="replace")[:max_chars]
                return {"success": True, "text": text, "pages": 1, "tables": [], "metadata": {}}
            else:
                return {"success": False, "error": f"Unsupported file type: {ext}"}
        except Exception as e:
            return {"success": False, "error": str(e)}

    @staticmethod
    def _read_pdf(path: str, max_chars: int) -> Dict:
        if not _PDF_OK:
            return {"success": False, "error": "pdfplumber not installed — pip install pdfplumber"}

        pages_text = []
        tables     = []
        metadata   = {}

        with pdfplumber.open(path) as pdf:
            metadata = pdf.metadata or {}
            for i, page in enumerate(pdf.pages):
                text = page.extract_text() or ""
                pages_text.append(text)
                for tbl in page.extract_tables():
                    if tbl:
                        tables.append({"page": i + 1, "rows": tbl})

        full_text = "\n\n".join(pages_text)
        if len(full_text) > max_chars:
            full_text = full_text[:max_chars] + f"\n... [{len(full_text)-max_chars} chars truncated]"

        return {
            "success":  True,
            "text":     full_text,
            "pages":    len(pages_text),
            "tables":   tables[:20],  # cap table count
            "metadata": {k: str(v) for k, v in metadata.items() if v},
        }

    @staticmethod
    def _read_docx(path: str, max_chars: int) -> Dict:
        if not _DOCX_OK:
            return {"success": False, "error": "python-docx not installed — pip install python-docx"}

        doc   = _DocxDocument(path)
        paras = [p.text for p in doc.paragraphs if p.text.strip()]
        tables = []
        for tbl in doc.tables:
            rows = []
            for row in tbl.rows:
                rows.append([cell.text.strip() for cell in row.cells])
            if rows:
                tables.append({"rows": rows})

        text = "\n".join(paras)
        if len(text) > max_chars:
            text = text[:max_chars] + f"\n... truncated"

        return {"success": True, "text": text, "pages": len(paras), "tables": tables, "metadata": {}}

    @staticmethod
    def _read_pptx(path: str, max_chars: int) -> Dict:
        if not _PPTX_OK:
            return {"success": False, "error": "python-pptx not installed — pip install python-pptx"}

        prs    = _Pptx(path)
        slides = []
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            if slide_text:
                slides.append(f"[Slide {i+1}]\n" + "\n".join(slide_text))

        text = "\n\n".join(slides)
        if len(text) > max_chars:
            text = text[:max_chars] + "\n... truncated"

        return {
            "success": True,
            "text":    text,
            "pages":   len(prs.slides),
            "tables":  [],
            "metadata": {"slide_count": len(prs.slides)},
        }

    # ── Image analysis ─────────────────────────────────────────────────────────

    @staticmethod
    def analyze_image(path: str, question: str = "") -> Dict:
        """Send an image to Claude vision and return the analysis."""
        if not Path(path).exists():
            return {"success": False, "error": f"File not found: {path}"}

        ext = Path(path).suffix.lower()
        if ext not in (".jpg", ".jpeg", ".png", ".gif", ".webp", ".bmp", ".tiff"):
            return {"success": False, "error": f"Unsupported image format: {ext}"}

        # Convert non-native formats to JPEG first
        actual_path = path
        if ext in (".bmp", ".tiff") and _PIL_OK:
            tmp = str(_TEMP_DIR / f"conv_{Path(path).stem}.jpg")
            _ensure_dir(_TEMP_DIR)
            _PIL.open(path).convert("RGB").save(tmp, "JPEG")
            actual_path = tmp

        q = question or "Describe this image in detail."
        try:
            # T-121: provider chain — gemini_flash → gemini_pro → claude_haiku → claude_sonnet
            # 429s parsed for retryDelay and advanced past; 3 consecutive fails open 5-min circuit.
            result, backend = _vision_via_router([actual_path], q)
            return {"success": True, "path": path, "analysis": result, "backend": backend}
        except Exception as e:
            return {"success": False, "error": str(e)}

    @staticmethod
    def analyze_images(paths: List[str], question: str = "") -> Dict:
        """Analyze multiple images together (compare, find differences, etc.)."""
        existing = [p for p in paths if Path(p).exists()]
        if not existing:
            return {"success": False, "error": "No valid image paths provided"}
        try:
            result = _claude_vision(existing, question or "Analyze these images.")
            return {"success": True, "paths": existing, "analysis": result}
        except Exception as e:
            return {"success": False, "error": str(e)}

    # ── Video analysis ─────────────────────────────────────────────────────────

    @staticmethod
    def analyze_video(path: str, question: str = "", max_frames: int = 8) -> Dict:
        """
        Sample frames from a video and analyze with Claude vision.
        Samples up to max_frames evenly spaced frames.
        """
        if not _CV2_OK:
            return {"success": False, "error": "opencv-python not installed — pip install opencv-python"}
        if not Path(path).exists():
            return {"success": False, "error": f"File not found: {path}"}

        _ensure_dir(_TEMP_DIR)
        cap = _cv2.VideoCapture(path)
        if not cap.isOpened():
            return {"success": False, "error": "Could not open video file"}

        total_frames = int(cap.get(_cv2.CAP_PROP_FRAME_COUNT))
        fps          = cap.get(_cv2.CAP_PROP_FPS)
        duration_s   = total_frames / fps if fps > 0 else 0

        if total_frames < 1:
            cap.release()
            return {"success": False, "error": "Video has no frames"}

        # Sample evenly
        frame_indices = [
            int(i * total_frames / max_frames)
            for i in range(min(max_frames, total_frames))
        ]

        saved_frames = []
        for idx in frame_indices:
            cap.set(_cv2.CAP_PROP_POS_FRAMES, idx)
            ret, frame = cap.read()
            if not ret:
                continue
            frame_path = str(_TEMP_DIR / f"frame_{idx:06d}.jpg")
            _cv2.imwrite(frame_path, frame)
            saved_frames.append(frame_path)

        cap.release()

        if not saved_frames:
            return {"success": False, "error": "Could not extract frames"}

        try:
            q = question or f"Analyze this video ({duration_s:.1f}s). Describe what's happening across all frames."
            # T-121: provider chain — same router as analyze_image
            analysis, backend = _vision_via_router(saved_frames, q)
            return {
                "success":      True,
                "path":         path,
                "duration_s":   round(duration_s, 1),
                "total_frames": total_frames,
                "sampled":      len(saved_frames),
                "analysis":     analysis,
                "backend":      backend,
            }
        except Exception as e:
            return {"success": False, "error": str(e)}
        finally:
            # Clean up temp frames
            for fp in saved_frames:
                try:
                    os.remove(fp)
                except Exception:
                    pass

    # ── OCR (T-052) ───────────────────────────────────────────────────────────

    @staticmethod
    def ocr_image(path: str, strategy: str = "auto") -> Dict:
        """Extract text from an image.

        Args:
            path: Path to image file.
            strategy: "auto" (Tesseract → Claude fallback), "tesseract" (only),
                      or "claude" (only). Default "auto".

        Returns:
            {"success", "text", "engine", "confidence"}
        """
        if not Path(path).exists():
            return {"success": False, "error": f"File not found: {path}"}

        # Tesseract path
        if strategy in ("auto", "tesseract") and _TESS_OK and _PIL_OK:
            try:
                img = _PIL.open(path)
                data = pytesseract.image_to_data(img, output_type=pytesseract.Output.DICT)
                # Compute mean confidence over non-empty words
                confs = [int(c) for c, w in zip(data["conf"], data["text"])
                         if w.strip() and c != "-1"]
                text = pytesseract.image_to_string(img).strip()
                if text:
                    confidence = round(sum(confs) / len(confs) / 100, 3) if confs else 0.0
                    if strategy == "auto" and confidence < 0.5:
                        pass  # low confidence → fall through to Claude
                    else:
                        return {
                            "success": True,
                            "text": text,
                            "engine": "tesseract",
                            "confidence": confidence,
                        }
            except Exception:
                pass

        if strategy == "tesseract":
            return {"success": False, "error": "Tesseract OCR failed or not installed"}

        # Claude vision path
        try:
            result = _claude_vision(
                [path],
                "Extract ALL text from this image exactly as it appears. "
                "Return only the extracted text, no commentary or formatting.",
            )
            return {
                "success": True,
                "text": result,
                "engine": "claude_vision",
                "confidence": 1.0,
            }
        except Exception as e:
            return {"success": False, "error": str(e)}

    # ── Facial recognition ────────────────────────────────────────────────────

    @staticmethod
    def detect_faces(path: str) -> Dict:
        """
        Detect faces in an image and return attributes (age, gender, emotion, race).
        Uses DeepFace. Does NOT identify who the person is — use recognize_face for that.
        """
        if _get_deepface() is None:
            return {"success": False, "error": "deepface not installed — pip install deepface"}
        if not Path(path).exists():
            return {"success": False, "error": f"File not found: {path}"}

        try:
            results = _get_deepface().analyze(
                img_path=path,
                actions=["age", "gender", "emotion", "race"],
                enforce_detection=False,
                silent=True,
            )
            if not isinstance(results, list):
                results = [results]

            faces = []
            for r in results:
                faces.append({
                    "age":           r.get("age"),
                    "gender":        r.get("dominant_gender"),
                    "emotion":       r.get("dominant_emotion"),
                    "race":          r.get("dominant_race"),
                    "region":        r.get("region"),
                    "gender_scores": r.get("gender"),
                    "emotion_scores": r.get("emotion"),
                })

            return {
                "success":    True,
                "face_count": len(faces),
                "faces":      faces,
            }
        except Exception as e:
            return {"success": False, "error": str(e)}

    @staticmethod
    def register_face(path: str, name: str) -> Dict:
        """
        Register a face in the local database under the given name.
        Saves to data/faces/<name>/ directory (gitignored).
        """
        if not Path(path).exists():
            return {"success": False, "error": f"File not found: {path}"}
        if not name.strip():
            return {"success": False, "error": "Name cannot be empty"}

        safe_name = re.sub(r"[^\w\-]", "_", name.strip())
        face_dir  = _FACES_DIR / safe_name
        face_dir.mkdir(parents=True, exist_ok=True)

        # Count existing images for this person
        existing = list(face_dir.glob("*.jpg")) + list(face_dir.glob("*.png"))
        dest = face_dir / f"face_{len(existing)+1:03d}{Path(path).suffix}"

        try:
            shutil.copy2(path, dest)
            return {
                "success": True,
                "name":    safe_name,
                "saved":   str(dest),
                "total_for_person": len(existing) + 1,
            }
        except Exception as e:
            return {"success": False, "error": str(e)}

    @staticmethod
    def recognize_face(path: str, threshold: float = 0.4) -> Dict:
        """
        Identify a face in path against the registered face database.
        Returns the best match name + confidence. Needs registered faces via register_face().
        """
        if _get_deepface() is None:
            return {"success": False, "error": "deepface not installed — pip install deepface"}
        if not Path(path).exists():
            return {"success": False, "error": f"File not found: {path}"}
        if not _FACES_DIR.exists() or not any(_FACES_DIR.iterdir()):
            return {"success": False, "error": "No faces registered. Use register_face() first."}

        try:
            results = _get_deepface().find(
                img_path=path,
                db_path=str(_FACES_DIR),
                model_name="VGG-Face",
                enforce_detection=False,
                silent=True,
            )
            matches = []
            for df in (results if isinstance(results, list) else [results]):
                if hasattr(df, "iterrows"):
                    for _, row in df.iterrows():
                        identity = str(row.get("identity", ""))
                        # Extract name from path: data/faces/<name>/face_001.jpg
                        parts    = Path(identity).parts
                        try:
                            name_idx = parts.index("faces") + 1
                            name     = parts[name_idx]
                        except (ValueError, IndexError):
                            name = Path(identity).parent.name
                        dist = float(row.get("distance", 1.0))
                        matches.append({"name": name, "distance": round(dist, 4), "path": identity})

            if not matches:
                return {"success": True, "match": None, "message": "No match found in database"}

            best = min(matches, key=lambda x: x["distance"])
            confident = best["distance"] <= threshold

            return {
                "success":    True,
                "match":      best["name"] if confident else None,
                "distance":   best["distance"],
                "confident":  confident,
                "threshold":  threshold,
                "all_matches": matches[:5],
            }
        except Exception as e:
            return {"success": False, "error": str(e)}

    @staticmethod
    def list_registered_faces() -> Dict:
        """List all registered people in the face database."""
        if not _FACES_DIR.exists():
            return {"success": True, "people": [], "count": 0}
        people = []
        for d in sorted(_FACES_DIR.iterdir()):
            if d.is_dir():
                imgs = len(list(d.glob("*.jpg")) + list(d.glob("*.png")))
                people.append({"name": d.name, "images": imgs})
        return {"success": True, "people": people, "count": len(people)}

    # ── Smart document analysis (T-051) ──────────────────────────────────────

    @staticmethod
    def analyze_document_smart(path: str, question: str = "", max_vision_pages: int = 5) -> Dict:
        """Two-pass document analysis for PDFs with charts, figures, or mixed content.

        Pass 1 — text density scan: score every page by character count and visual
        indicators (low text + image placeholders = likely chart/figure page).
        Pass 2 — vision on top-N: render the highest-value pages as images and
        send them to Claude vision for visual analysis.

        Falls back to read_document for non-PDF files.
        """
        p = Path(path)
        if not p.exists():
            return {"success": False, "error": f"File not found: {path}"}

        ext = p.suffix.lower()
        if ext != ".pdf":
            # Non-PDF: just text extraction
            result = MediaTools.read_document(path)
            result["vision_pages"] = []
            result["vision_analysis"] = ""
            return result

        if not _PDF_OK:
            return {"success": False, "error": "pdfplumber not installed — pip install pdfplumber"}

        try:
            import pdfplumber

            page_scores: List[Tuple[int, float]] = []  # (page_index, score)
            pages_text: List[str] = []
            tables: List[Dict] = []
            metadata: Dict = {}

            with pdfplumber.open(path) as pdf:
                metadata = pdf.metadata or {}
                for i, page in enumerate(pdf.pages):
                    text = page.extract_text() or ""
                    pages_text.append(text)

                    # Score = visual richness (inverse of text density)
                    char_count = len(text.strip())
                    img_count = len(page.images) if hasattr(page, "images") else 0
                    # Low text + images = chart/figure = high score for vision
                    score = img_count * 30 + max(0, 500 - char_count)
                    page_scores.append((i, score))

                    for tbl in page.extract_tables():
                        if tbl:
                            tables.append({"page": i + 1, "rows": tbl})

            # Select top-N pages for vision analysis (by score, skip purely blank pages)
            top_pages = sorted(page_scores, key=lambda x: x[1], reverse=True)[:max_vision_pages]
            top_page_indices = sorted(idx for idx, _ in top_pages if pages_text[idx].strip() or _)

            full_text = "\n\n".join(pages_text)
            if len(full_text) > 50_000:
                full_text = full_text[:50_000] + f"\n... [{len(full_text)-50_000} chars truncated]"

            vision_analysis = ""
            vision_pages_used: List[int] = []

            if _PIL_OK and top_page_indices:
                _ensure_dir(_TEMP_DIR)
                tmp_images: List[str] = []
                try:
                    with pdfplumber.open(path) as pdf:
                        for idx in top_page_indices:
                            try:
                                pg_img = pdf.pages[idx].to_image(resolution=150)
                                tmp_path = str(_TEMP_DIR / f"smart_p{idx:04d}.png")
                                pg_img.save(tmp_path)
                                tmp_images.append(tmp_path)
                                vision_pages_used.append(idx + 1)  # 1-based
                            except Exception:
                                pass

                    if tmp_images:
                        q = question or (
                            f"Analyze these {len(tmp_images)} pages from a PDF. "
                            "Extract all text visible in charts, diagrams, figures, tables, "
                            "and annotate what each visual element shows."
                        )
                        vision_analysis = _claude_vision(tmp_images, q, max_tokens=4096)
                finally:
                    for f in tmp_images:
                        try:
                            os.remove(f)
                        except Exception:
                            pass

            combined = ""
            if full_text and vision_analysis:
                combined = (
                    f"=== EXTRACTED TEXT ({len(pages_text)} pages) ===\n{full_text}\n\n"
                    f"=== VISION ANALYSIS (pages {vision_pages_used}) ===\n{vision_analysis}"
                )
            elif full_text:
                combined = full_text
            else:
                combined = vision_analysis

            return {
                "success": True,
                "path": path,
                "pages": len(pages_text),
                "text": full_text,
                "tables": tables[:20],
                "vision_pages": vision_pages_used,
                "vision_analysis": vision_analysis,
                "combined": combined[:20_000],
                "metadata": {k: str(v) for k, v in metadata.items() if v},
            }

        except Exception as e:
            return {"success": False, "error": str(e)}

    # ── Document + image combo (legacy) ──────────────────────────────────────

    @staticmethod
    def analyze_document_with_vision(path: str, question: str = "") -> Dict:
        """
        Smart document analysis: extracts text for text-based docs,
        uses Claude vision for scanned/image-heavy PDFs.
        Returns {"text", "vision_analysis", "combined"}.
        """
        ext = Path(path).suffix.lower()

        text_result = MediaTools.read_document(path)
        text        = text_result.get("text", "") if text_result.get("success") else ""

        vision_analysis = ""
        if ext == ".pdf" and _PIL_OK and _PDF_OK:
            # Also render first page as image for visual analysis
            try:
                import pdfplumber
                with pdfplumber.open(path) as pdf:
                    if pdf.pages:
                        pg_img = pdf.pages[0].to_image(resolution=150)
                        _ensure_dir(_TEMP_DIR)
                        tmp_img = str(_TEMP_DIR / "page1.png")
                        pg_img.save(tmp_img)
                        q = question or "Analyze this document page — extract any visual elements, diagrams, charts, signatures, or important layout information not captured in plain text."
                        vision_analysis = _claude_vision([tmp_img], q)
                        try:
                            os.remove(tmp_img)
                        except Exception:
                            pass
            except Exception as e:
                track_silent("media.analyze_document_with_vision_pdf", e)

        elif ext in (".png", ".jpg", ".jpeg", ".gif", ".webp"):
            r = MediaTools.analyze_image(path, question)
            vision_analysis = r.get("analysis", "")

        combined = ""
        if text and vision_analysis:
            combined = f"=== EXTRACTED TEXT ===\n{text}\n\n=== VISUAL ANALYSIS ===\n{vision_analysis}"
        elif text:
            combined = text
        elif vision_analysis:
            combined = vision_analysis

        return {
            "success":         True,
            "path":            path,
            "text":            text[:10000] if text else "",
            "vision_analysis": vision_analysis,
            "combined":        combined[:15000],
            "pages":           text_result.get("pages", 0),
            "tables":          text_result.get("tables", []),
        }


# ── T-083 R2.1: tool registry export ─────────────────────────────────────────
#
# MediaTools static methods (no instance state) — handlers call the class
# directly rather than instantiating, matching the legacy _media (lazy class
# proxy) pattern.

from agent.tool_spec import ToolSpec  # noqa: E402


def _handle_read_document(agent, tool_input, *, memory_override=None):
    return MediaTools.read_document(
        path=tool_input["path"],
        max_chars=tool_input.get("max_chars", 50_000),
    )


def _handle_analyze_image(agent, tool_input, *, memory_override=None):
    return MediaTools.analyze_image(
        path=tool_input["path"],
        question=tool_input.get("question", ""),
    )


def _handle_analyze_images(agent, tool_input, *, memory_override=None):
    return MediaTools.analyze_images(
        paths=tool_input["paths"],
        question=tool_input.get("question", ""),
    )


def _handle_analyze_video(agent, tool_input, *, memory_override=None):
    return MediaTools.analyze_video(
        path=tool_input["path"],
        question=tool_input.get("question", ""),
        max_frames=tool_input.get("max_frames", 8),
    )


def _handle_ocr_image(agent, tool_input, *, memory_override=None):
    return MediaTools.ocr_image(path=tool_input["path"])


def _handle_analyze_document_smart(agent, tool_input, *, memory_override=None):
    return MediaTools.analyze_document_with_vision(
        path=tool_input["path"],
        question=tool_input.get("question", ""),
    )


def _handle_analyze_media(agent, tool_input, *, memory_override=None):
    """Merged handler — routes to image/video/ocr/document based on kind or extension."""
    raw = tool_input.get("paths") or tool_input.get("path")
    if isinstance(raw, str):
        paths = [raw]
    elif isinstance(raw, list):
        paths = list(raw)
    else:
        return {"success": False, "error": "Provide 'path' (string) or 'paths' (list)."}

    question = tool_input.get("question", "")
    kind = tool_input.get("kind", "auto")

    if kind == "auto":
        ext = Path(paths[0]).suffix.lower() if paths else ""
        if ext in (".mp4", ".avi", ".mov", ".mkv", ".webm", ".flv", ".wmv"):
            kind = "video"
        elif ext in (".jpg", ".jpeg", ".png", ".gif", ".webp", ".bmp", ".tiff"):
            kind = "image"
        else:
            kind = "document"

    if kind == "video":
        return MediaTools.analyze_video(
            path=paths[0], question=question,
            max_frames=tool_input.get("max_frames", 8),
        )
    if kind == "ocr":
        return MediaTools.ocr_image(
            path=paths[0], strategy=tool_input.get("strategy", "auto"),
        )
    if kind == "document":
        return MediaTools.analyze_document_smart(
            path=paths[0], question=question,
        )
    # kind == "image" (default)
    if len(paths) == 1:
        return MediaTools.analyze_image(path=paths[0], question=question)
    return MediaTools.analyze_images(paths=paths, question=question)


def _handle_detect_faces(agent, tool_input, *, memory_override=None):
    return MediaTools.detect_faces(path=tool_input["path"])


def _handle_recognize_face(agent, tool_input, *, memory_override=None):
    return MediaTools.recognize_face(
        path=tool_input["path"],
        threshold=tool_input.get("threshold", 0.4),
    )


def _handle_register_face(agent, tool_input, *, memory_override=None):
    return MediaTools.register_face(
        path=tool_input["path"],
        name=tool_input["name"],
    )


def _handle_list_registered_faces(agent, tool_input, *, memory_override=None):
    return MediaTools.list_registered_faces()


TOOLS = [
    ToolSpec(
        name="read_document",
        description="Extract text and tables from PDF, DOCX, PPTX, or plain text files. Returns full text content.",
        input_schema={
            "type": "object",
            "properties": {
                "path":      {"type": "string", "description": "Absolute or relative file path"},
                "max_chars": {"type": "integer", "default": 50000},
            },
            "required": ["path"],
        },
        handler=_handle_read_document,
        success_predicate=lambda r: r.get("success", False),
    ),
    ToolSpec(
        name="analyze_media",
        description=(
            "Analyze any media file — image(s), video, scanned PDF, or document. "
            "Auto-detects the right analysis from the file extension. "
            "Pass kind='image'/'video'/'ocr'/'document' to override. "
            "Supports JPG, PNG, GIF, WebP, MP4, AVI, MOV, PDF, DOCX, PPTX, TXT."
        ),
        input_schema={
            "type": "object",
            "properties": {
                "path":       {"type": "string",
                               "description": "Single file path (image, video, or document)"},
                "paths":      {"type": "array", "items": {"type": "string"},
                               "description": "Multiple image paths for side-by-side analysis"},
                "question":   {"type": "string",
                               "description": "What to analyze, extract, or ask about the file"},
                "kind":       {"type": "string",
                               "enum": ["auto", "image", "video", "ocr", "document"],
                               "default": "auto",
                               "description": "Analysis type — 'auto' detects from extension"},
                "max_frames": {"type": "integer", "default": 8,
                               "description": "Video: how many frames to sample"},
            },
            "required": [],
        },
        handler=_handle_analyze_media,
        success_predicate=lambda r: r.get("success", False),
        aliases=("analyze_image", "analyze_images", "analyze_video", "ocr_image", "analyze_document_smart"),
    ),
    ToolSpec(
        name="detect_faces",
        description="Detect faces in an image. Returns count + bounding boxes. Uses MediaPipe (CPU, no GPU needed).",
        input_schema={
            "type": "object",
            "properties": {"path": {"type": "string"}},
            "required": ["path"],
        },
        handler=_handle_detect_faces,
        success_predicate=lambda r: r.get("success", False),
    ),
    ToolSpec(
        name="recognize_face",
        description="Recognize known faces in an image. Returns matched names + confidence. Requires register_face first.",
        input_schema={
            "type": "object",
            "properties": {
                "path":      {"type": "string"},
                "threshold": {"type": "number", "default": 0.4,
                              "description": "Confidence threshold (0-1)"},
            },
            "required": ["path"],
        },
        handler=_handle_recognize_face,
        success_predicate=lambda r: r.get("success", False),
    ),
    ToolSpec(
        name="register_face",
        description="Register a person's face for later recognition. Pass a clear photo and their name.",
        input_schema={
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "Clear photo of the person"},
                "name": {"type": "string", "description": "Person's name"},
            },
            "required": ["path", "name"],
        },
        handler=_handle_register_face,
        success_predicate=lambda r: r.get("success", False),
    ),
    ToolSpec(
        name="list_registered_faces",
        description="List all registered faces with metadata.",
        input_schema={"type": "object", "properties": {}, "required": []},
        handler=_handle_list_registered_faces,
        success_predicate=lambda r: r.get("success", False),
    ),
]
